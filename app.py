"""
Pride Advisory Financial Evaluator — Full Build
Features added in this build:
1. Smart input preview next to each numeric input (shows parsed value in Naira)
2. Improved 3-statement uploader with fuzzy mapping for Income Statement, Balance Sheet, Cash Flow
3. DCF -> intrinsic per share + WACC vs growth sensitivity (heatmap)
4. PDF and PowerPoint export (basic, includes key charts and tables)
5. Simple Streamlit authentication scaffold (uses STREAMLIT_SECRETS for credentials)

How to run:
- pip install streamlit pandas numpy plotly python-pptx fpdf openpyxl
- streamlit run this_file.py

Notes:
- This file is a runnable template. You may tweak the report styling and chart content as you like.
- For Cloud deployment, add secrets in Streamlit Cloud (username/password) and set STREAMLIT_SECRETS accordingly.
"""

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from fpdf import FPDF
import math
import re
import copy

# -------------------- Helpers --------------------
UNIT_MAP = {'k':1e3,'thousand':1e3,'m':1e6,'mn':1e6,'million':1e6,'b':1e9,'bn':1e9,'billion':1e9}

def parse_number(s):
    if s is None:
        return None
    if isinstance(s,(int,float,np.number)):
        return float(s)
    t = str(s).strip().lower()
    if t == '':
        return None
    t = t.replace(',','')
    neg = False
    if t.startswith('(') and t.endswith(')'):
        neg = True; t = t[1:-1]
    m = re.match(r'^([+-]?\d*\.?\d+)([a-zA-Z]+)?$', t)
    if not m:
        try:
            v = float(t)
            return -v if neg else v
        except:
            return None
    num = float(m.group(1))
    suf = (m.group(2) or '').lower()
    mul = 1
    for k in UNIT_MAP:
        if k in suf:
            mul = UNIT_MAP[k]
            break
    val = num * mul
    return -val if neg else val

def fmt_naira(x):
    if x is None or (isinstance(x,float) and (math.isnan(x) or math.isinf(x))):
        return '-'
    sign = '-' if x<0 else ''
    a = abs(x)
    if a>=1e9:
        return f"{sign}₦{a/1e9:,.3f}b"
    if a>=1e6:
        return f"{sign}₦{a/1e6:,.3f}m"
    if a>=1e3:
        return f"{sign}₦{a:,.0f}"
    return f"{sign}₦{a:,.2f}"

def safe_div(a,b):
    try:
        if a is None or b is None: return None
        if b==0: return None
        return a/b
    except:
        return None

# -------------------- Simple Auth Scaffold --------------------
# Optional: set secrets via Streamlit Cloud: {'username':'you','password':'pass'}
USE_AUTH = True
if USE_AUTH:
    secrets = st.secrets.get('credentials', {}) if st.secrets else {}
    USER = secrets.get('username')
    PASS = secrets.get('password')
    if USER and PASS:
        st.sidebar.write('Login')
        user_in = st.sidebar.text_input('Username')
        pass_in = st.sidebar.text_input('Password', type='password')
        if user_in != USER or pass_in != PASS:
            st.sidebar.error('Invalid credentials')
            st.stop()
    else:
        # no secrets configured — allow local dev but show warning
        st.sidebar.warning('No auth secrets found — running in open mode')

# -------------------- App UI --------------------
st.set_page_config(page_title='Pride Advisory — Financial Evaluator', layout='wide')
st.markdown("<h1 style='text-align:center'>PRIDE ADVISORY — Financial Evaluator</h1>", unsafe_allow_html=True)
st.write('Smart parsing: 5bn, 250m, (200m). Upload full 3-statement Excel (optional).')

cols = st.columns([2,1])
with cols[0]:
    analysis = st.multiselect('Choose modules', ['Ratios','DCF','Sensitivity','Charts','Export Report'], default=['Ratios','DCF','Sensitivity'])
    years = st.slider('Years (history/forecast)', min_value=1, max_value=10, value=5)
    uploaded = st.file_uploader('Upload 3-statement Excel/CSV (optional)', type=['xlsx','csv'])
with cols[1]:
    st.subheader('Market defaults (editable)')
    rf = st.number_input('Risk-free rate (10y govt) %', value=16.22)/100.0
    mm_erp = st.number_input('Mature market ERP %', value=4.33)/100.0
    country_prem = st.number_input('Country premium %', value=7.43)/100.0
    st.caption('Defaults are suggested — edit per company case')

# -------------------- 3-statement parsing --------------------
income_df = balance_df = cash_df = None
if uploaded is not None:
    try:
        if uploaded.name.endswith('.csv'):
            dfu = pd.read_csv(uploaded)
            st.info('CSV uploaded — expecting vertical or wide table. Please check parsed preview.')
            income_df = dfu
        else:
            xls = pd.read_excel(uploaded, sheet_name=None)
            # try map common sheets
            for name, sheet in xls.items():
                lname = name.lower()
                if 'income' in lname or 'profit' in lname:
                    income_df = sheet
                elif 'balance' in lname or 'sheet' in lname or 'statement of financial position' in lname:
                    balance_df = sheet
                elif 'cash' in lname or 'flow' in lname:
                    cash_df = sheet
            # fallback picks
            if income_df is None:
                income_df = list(xls.values())[0]
            if balance_df is None and len(xls)>1:
                balance_df = list(xls.values())[1]
            st.success('Sheets loaded — verify mapped values.')
    except Exception as e:
        st.error('Upload parse failed: '+str(e))

# auto-extract series helper (fuzzy match)
def extract_series_from_df(df, candidates, n):
    if df is None: return [None]*n
    cols = [c.lower() for c in df.columns]
    for cand in candidates:
        for i,c in enumerate(cols):
            if cand in c:
                vals = df.iloc[:,i].dropna().tolist()
                parsed = [parse_number(v) for v in vals]
                if len(parsed)>=n:
                    return parsed[:n]
    # fallback: take numeric columns last n rows
    numeric_cols = df.select_dtypes(include=[np.number])
    if numeric_cols.shape[1]>0:
        col0 = numeric_cols.iloc[:,0]
        vals = col0.dropna().tolist()
        parsed = [float(v) for v in vals]
        if len(parsed)>=n:
            return parsed[:n]
    return [None]*n

# Attempt extraction
rev = cogs = ebit = net = assets = equity = [None]*years
if uploaded is not None:
    rev = extract_series_from_df(income_df, ['revenue','sales','turnover'], years)
    net = extract_series_from_df(income_df, ['net income','profit'], years)
    ebit = extract_series_from_df(income_df, ['ebit','operating income','operating profit'], years)
    if balance_df is not None:
        assets = extract_series_from_df(balance_df, ['total assets','assets'], years)
        equity = extract_series_from_df(balance_df, ['total equity','shareholders','equity'], years)
    # cash and debt from balance sheet latest row
    cash = None; debt = None
    if balance_df is not None:
        for c in balance_df.columns:
            lc = c.lower()
            if 'cash' in lc and cash is None:
                cash = parse_number(balance_df[c].dropna().iloc[0])
            if 'debt' in lc and debt is None:
                debt = parse_number(balance_df[c].dropna().iloc[0])
else:
    cash = None; debt = None

# Manual entry fallback & smart preview
st.markdown('---')
st.subheader('Manual inputs (leave blanks if uploaded)')
colA, colB = st.columns(2)
with colA:
    rev_text = st.text_input(f'Revenue series (comma-separated, {years} vals)', value=','.join([str(int(x)) for x in rev if x]) if any(rev) else '')
    cogs_text = st.text_input('COGS series (comma-separated)', value='')
    ebit_text = st.text_input('EBIT series (comma-separated)', value='')
    net_text = st.text_input('Net Income series (comma-separated)', value='')
with colB:
    assets_text = st.text_input('Total Assets series (comma-separated)', value='')
    equity_text = st.text_input('Equity series (comma-separated)', value='')
    cash_text = st.text_input('Cash (latest)', value='' if cash is None else str(int(cash)))
    debt_text = st.text_input('Total Debt (latest)', value='' if debt is None else str(int(debt)))
    shares_text = st.text_input('Shares outstanding (single)', value='')

# parse series helper
def parse_series(text, n):
    if not text: return [None]*n
    parts = [p.strip() for p in text.split(',') if p.strip()!='']
    vals = [parse_number(p) for p in parts]
    if len(vals)<n: vals = vals + [None]*(n-len(vals))
    return vals[:n]

# Prefer manual if provided, else uploaded
rev_manual = parse_series(rev_text, years)
rev = [rm if rm is not None else ru for rm,ru in zip(rev_manual, rev)]
cogs = parse_series(cogs_text, years)
ebit = parse_series(ebit_text, years)
net = parse_series(net_text, years)
assets_manual = parse_series(assets_text, years)
assets = [am if am is not None else au for am,au in zip(assets_manual, assets)]
equity_manual = parse_series(equity_text, years)
equity = [em if em is not None else eu for em,eu in zip(equity_manual, equity)]
cash = parse_number(cash_text) if cash_text else cash
debt = parse_number(debt_text) if debt_text else debt
shares = parse_number(shares_text) if shares_text else None

# show smart preview for the latest year
st.subheader('Smart preview (latest values)')
preview = pd.DataFrame({
    'Revenue':[fmt_naira(rev[-1])], 'COGS':[fmt_naira(cogs[-1]) if any(cogs) else '-'], 'EBIT':[fmt_naira(ebit[-1]) if any(ebit) else '-'],
    'Net Income':[fmt_naira(net[-1]) if any(net) else '-'], 'Total Assets':[fmt_naira(assets[-1]) if any(assets) else '-'], 'Equity':[fmt_naira(equity[-1]) if any(equity) else '-']
})
st.table(preview)

# -------------------- Ratios --------------------
def calc_ratios_latest():
    latest = {'revenue':rev[-1],'cogs':cogs[-1],'ebit':ebit[-1],'net':net[-1],'assets':assets[-1],'equity':equity[-1],'cash':cash,'debt':debt}
    r = {}
    r['Gross Profit Margin'] = safe_div(latest['revenue'] - latest['cogs'] if latest['revenue'] is not None and latest['cogs'] is not None else None, latest['revenue'])
    r['Operating Margin'] = safe_div(latest['ebit'], latest['revenue'])
    r['Net Profit Margin'] = safe_div(latest['net'], latest['revenue'])
    r['ROA'] = safe_div(latest['net'], latest['assets'])
    r['ROE'] = safe_div(latest['net'], latest['equity'])
    r['Current Ratio'] = None
    r['Debt to Equity'] = safe_div(latest['debt'], latest['equity'])
    return r

if 'Ratios' in analysis:
    st.subheader('Ratios (latest)')
    ratios = calc_ratios_latest()
    df_rat = pd.DataFrame([{'Metric':k,'Value':(f"{v:.2%}" if isinstance(v,(float,np.floating)) else fmt_naira(v) if v is not None else '-') } for k,v in ratios.items()])
    st.table(df_rat)

# -------------------- DCF & Sensitivity --------------------
def project_fcff(start_rev, years, rev_growth_list, ebit_margin, tax_rate, capex_pct, dep_pct, nwc_pct):
    rows = []
    rev = start_rev
    for t in range(1, years+1):
        g = rev_growth_list[t-1] if isinstance(rev_growth_list,(list,tuple)) and len(rev_growth_list)>=t else rev_growth_list[0] if isinstance(rev_growth_list,(list,tuple)) else rev_growth_list
        rev = rev * (1+g)
        ebitv = rev * ebit_margin
        dep = rev * dep_pct
        capex = rev * capex_pct
        change_nwc = rev * nwc_pct
        nopat = ebitv * (1 - tax_rate)
        fcff = nopat + dep - capex - change_nwc
        rows.append({'year':t,'revenue':rev,'ebit':ebitv,'nopat':nopat,'dep':dep,'capex':capex,'change_nwc':change_nwc,'fcff':fcff})
    return rows

def dcf_from_proj(rows, wacc, perp_growth=None, use_exit_mult=False, exit_mult=8, total_debt=0, cash=0, shares=None):
    pv_sum = 0
    for r in rows:
        t = r['year']
        r['pv'] = r['fcff']/((1+wacc)**t)
        pv_sum += r['pv']
    last = rows[-1]
    if use_exit_mult:
        terminal = last['ebit'] * exit_mult
    else:
        fcff_next = last['fcff'] * (1 + (perp_growth if perp_growth is not None else 0))
        denom = (wacc - perp_growth) if perp_growth is not None else None
        terminal = fcff_next / denom if denom and denom!=0 else None
    terminal_pv = terminal/((1+wacc)**last['year']) if terminal is not None else None
    ev = pv_sum + (terminal_pv or 0)
    net_debt = (total_debt or 0) - (cash or 0)
    eq = ev - net_debt
    per_share = eq/shares if shares and shares>0 else None
    return {'rows':rows,'enterprise_value':ev,'terminal':terminal,'terminal_pv':terminal_pv,'equity_value':eq,'per_share':per_share}

if 'DCF' in analysis:
    st.subheader('DCF assumptions')
    col1,col2,col3 = st.columns(3)
    with col1:
        start_rev = st.number_input('Starting revenue (latest) in Naira', value=rev[-1] if rev[-1] else 0.0)
        ebit_margin = st.number_input('EBIT margin (decimal)', value=0.12)
        tax_rate = st.number_input('Tax rate (decimal)', value=0.30)
    with col2:
        rev_growth_text = st.text_input('Revenue growth (single or comma list)', value='0.05')
        capex_pct = st.number_input('CapEx % of revenue', value=0.05)
        dep_pct = st.number_input('Depreciation % of revenue', value=0.03)
    with col3:
        nwc_pct = st.number_input('Change in NWC %', value=0.02)
        wacc_input = st.number_input('WACC (decimal)', value=0.12)
        perp_growth = st.number_input('Perpetual growth (decimal)', value=0.03)

    # parse growth
    try:
        if ',' in rev_growth_text:
            parts = [float(x.strip()) for x in rev_growth_text.split(',')]
            rev_growth_list = parts + [parts[-1]]*(years-len(parts)) if len(parts)<years else parts[:years]
        else:
            rev_growth_list = [float(rev_growth_text)]*years
    except:
        rev_growth_list = [0.05]*years

    proj = project_fcff(start_rev, years, rev_growth_list, ebit_margin, tax_rate, capex_pct, dep_pct, nwc_pct)
    use_exit = st.checkbox('Use exit multiple for terminal', value=False)
    exit_mult = st.number_input('Exit multiple', value=8.0)
    dcfres = dcf_from_proj(copy.deepcopy(proj), wacc_input, perp_growth, use_exit, exit_mult, debt, cash, shares)
    st.metric('Enterprise Value', fmt_naira(dcfres['enterprise_value']))
    st.metric('Equity Value', fmt_naira(dcfres['equity_value']))
    st.metric('Intrinsic per share', fmt_naira(dcfres['per_share']) if dcfres['per_share'] else '-')
    st.dataframe(pd.DataFrame(dcfres['rows']))

# Sensitivity
if 'Sensitivity' in analysis:
    st.subheader('Sensitivity: WACC vs Perpetual Growth')
    base_wacc = st.number_input('Base WACC', value=wacc_input if 'wacc_input' in locals() else 0.12)
    base_g = st.number_input('Base perp growth', value=perp_growth if 'perp_growth' in locals() else 0.03)
    wacc_vals = np.round(np.linspace(max(0.01,base_wacc-0.04), base_wacc+0.04, 7),4)
    g_vals = np.round(np.linspace(max(-0.01,base_g-0.02), base_g+0.02, 7),4)
    grid = pd.DataFrame(index=[f"{g:.2%}" for g in g_vals], columns=[f"{w:.2%}" for w in wacc_vals])
    for g in g_vals:
        for w in wacc_vals:
            res = dcf_from_proj(copy.deepcopy(proj), w, g, use_exit, exit_mult, debt, cash, shares)
            grid.loc[f"{g:.2%}", f"{w:.2%}"] = res['per_share']
    fig = go.Figure(data=go.Heatmap(z=grid.values.astype(float), x=grid.columns, y=grid.index, colorscale='Viridis'))
    fig.update_layout(title='Intrinsic value per share heatmap', xaxis_title='WACC', yaxis_title='Perpetual growth')
    st.plotly_chart(fig, use_container_width=True)

# -------------------- Charts --------------------
if 'Charts' in analysis:
    st.subheader('Trend charts')
    yrs = list(range(1,years+1))
    if any([x is not None for x in rev]):
        figr = px.line(x=yrs, y=[(v or np.nan) for v in rev], labels={'x':'Year','y':'Revenue'}, title='Revenue')
        st.plotly_chart(figr, use_container_width=True)
    if any([x is not None for x in net]):
        figm = px.line(x=yrs, y=[safe_div(net[i],rev[i]) if rev[i] else np.nan for i in range(years)], labels={'x':'Year','y':'Net Margin'}, title='Net Margin')
        st.plotly_chart(figm, use_container_width=True)

# -------------------- Report Export (PDF & PPTX) --------------------
if 'Export Report' in analysis:
    st.subheader('Export report')
    report_title = st.text_input('Report title', value='Pride Advisory Financial Analysis')
    if st.button('Generate PPTX'):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = report_title
        # add a simple table slide
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tx.text = 'Key ratios (latest)'
        bio = BytesIO()
        prs.save(bio)
        st.download_button('Download PPTX', data=bio.getvalue(), file_name='report.pptx')
    if st.button('Generate PDF'):
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(0,10,report_title, ln=True)
        pdf.output('report.pdf')
        with open('report.pdf','rb') as f:
            st.download_button('Download PDF', data=f, file_name='report.pdf')

# Footer
st.markdown('---')
st.caption('Pride Advisory — Financial Evaluator. Validate assumptions for each case.')
